"""
使用LibreOffice命令行触发PPT文本框自适应渲染
通过转换PPT为PDF的过程触发完整渲染，使文本框自适应设置真正生效
"""
import os
import logging
import platform
import subprocess
import tempfile
import time
import shutil
from typing import Optional, Dict, Any

try:
    from pptx import Presentation
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.util import Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.error("python-pptx不可用")

logger = logging.getLogger(__name__)


class LibreOfficeRenderTrigger:
    """LibreOffice渲染触发器"""

    def __init__(self):
        """初始化LibreOffice渲染触发器"""
        # 检查LibreOffice是否可用
        self.libreoffice_available = self._check_libreoffice()
        
        # 初始化统计信息
        self.stats = {
            'total_textboxes': 0,  # 总文本框数
            'autofit_set': 0,      # 设置了自适应的文本框数
            'render_triggered': 0, # 触发渲染次数
            'pdf_generated': 0,    # 生成的PDF文件数
            'pdf_deleted': 0,      # 删除的PDF文件数
            'pptx_generated': 0    # 生成的PPTX文件数
        }

    def _check_libreoffice(self) -> bool:
        """检查LibreOffice是否可用"""
        try:
            import platform

            # 根据操作系统定义可能的LibreOffice命令路径
            if platform.system() == "Windows":
                # Windows下的常见LibreOffice安装路径
                commands = [
                    r"C:\Program Files\LibreOffice\program\soffice.exe",
                    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                    r"soffice.exe",  # 如果在PATH中
                    r"libreoffice.exe"
                ]
            else:
                # Linux/macOS下的常见路径
                commands = [
                    'libreoffice',
                    'soffice',
                    '/usr/bin/libreoffice',
                    '/opt/libreoffice/program/soffice',
                    '/Applications/LibreOffice.app/Contents/MacOS/soffice'  # macOS
                ]

            for cmd in commands:
                try:
                    result = subprocess.run(
                        [cmd, '--version'],
                        capture_output=True,
                        text=True,
                        timeout=15
                    )
                    if result.returncode == 0:
                        self.libreoffice_cmd = cmd
                        version_info = result.stdout.strip()
                        logger.info(f"LibreOffice可用: {cmd}")
                        logger.info(f"版本信息: {version_info}")
                        return True
                except (FileNotFoundError, subprocess.TimeoutExpired, OSError):
                    continue

            logger.warning("LibreOffice不可用")
            logger.info("请确保LibreOffice已正确安装")
            return False

        except Exception as e:
            logger.warning(f"检查LibreOffice时出错: {e}")
            return False

    def process_ppt_with_render_trigger(self, ppt_path: str) -> bool:
        """
        处理PPT并通过LibreOffice触发渲染

        Args:
            ppt_path: PPT文件路径

        Returns:
            bool: 处理是否成功
        """
        if not os.path.exists(ppt_path):
            logger.error(f"文件不存在: {ppt_path}")
            return False

        if not self.libreoffice_available:
            logger.error("LibreOffice不可用，无法触发渲染")
            return False

        logger.info(f"开始处理PPT并触发渲染: {os.path.basename(ppt_path)}")

        try:
            # 步骤1: 使用python-pptx设置自适应属性
            if not self._set_autofit_properties(ppt_path):
                logger.error("设置自适应属性失败")
                return False

            # 步骤2: 使用LibreOffice转换PDF触发渲染
            if not self._trigger_render_via_pdf_conversion(ppt_path):
                logger.error("LibreOffice渲染触发失败")
                return False

            logger.info("✅ PPT文本框自适应处理完成（包含渲染触发）")
            self._log_stats()
            return True

        except Exception as e:
            logger.error(f"处理PPT时出错: {e}")
            return False

    def _set_autofit_properties(self, ppt_path: str) -> bool:
        """设置文本框自适应属性"""
        if not PPTX_AVAILABLE:
            logger.error("python-pptx库不可用")
            return False

        try:
            logger.info("步骤1: 设置文本框自适应属性")
            prs = Presentation(ppt_path)

            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        self.stats['total_textboxes'] += 1
                        text_frame = shape.text_frame

                        # 检查是否有实际内容
                        has_content = any(
                            paragraph.text.strip()
                            for paragraph in text_frame.paragraphs
                        )

                        if has_content:
                            # 设置文本大小适应文本框
                            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                            text_frame.word_wrap = True

                            # 优化边距
                            text_frame.margin_left = Pt(3)
                            text_frame.margin_right = Pt(3)
                            text_frame.margin_top = Pt(2)
                            text_frame.margin_bottom = Pt(2)

                            self.stats['autofit_set'] += 1

                    elif shape.has_table:
                        # 处理表格中的文本框
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                self.stats['total_textboxes'] += 1
                                text_frame = cell.text_frame

                                has_content = any(
                                    paragraph.text.strip()
                                    for paragraph in text_frame.paragraphs
                                )

                                if has_content:
                                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                                    text_frame.word_wrap = True
                                    self.stats['autofit_set'] += 1

            # 保存设置
            prs.save(ppt_path)
            logger.info(f"设置了 {self.stats['autofit_set']} 个文本框的自适应属性")
            return True

        except Exception as e:
            logger.error(f"设置自适应属性时出错: {e}")
            return False

    def _trigger_render_via_pdf_conversion(self, ppt_path: str) -> bool:
        """通过ODP中转转换为PDF触发渲染"""
        try:
            logger.info("步骤2: 通过LibreOffice ODP中转转换为PDF触发渲染")

            # 创建临时目录用于中间文件和PDF输出
            with tempfile.TemporaryDirectory() as temp_dir:
                # 确保路径格式正确
                temp_dir = os.path.abspath(temp_dir)
                ppt_path = os.path.abspath(ppt_path)
                ppt_filename = os.path.basename(ppt_path)
                ppt_name = os.path.splitext(ppt_filename)[0]

                logger.debug(f"输入文件: {ppt_path}")
                logger.debug(f"输出目录: {temp_dir}")

                # 步骤2.1: 先转换为ODP格式
                cmd_to_odp = [
                    self.libreoffice_cmd,
                    '--headless',           # 无头模式
                    '--invisible',          # 不可见
                    '--nodefault',          # 不使用默认设置
                    '--nolockcheck',        # 不检查锁定
                    '--nologo',             # 不显示logo
                    '--norestore',          # 不恢复
                    '--convert-to', 'odp',  # 转换为ODP
                    '--outdir', temp_dir,   # 输出目录
                    ppt_path                # 输入文件
                ]

                # 在Windows下，如果路径包含空格，需要特殊处理
                if platform.system() == "Windows":
                    # 对于Windows，确保路径被正确引用
                    cmd_str_odp = f'"{self.libreoffice_cmd}" --headless --invisible --nodefault --nolockcheck --nologo --norestore --convert-to odp --outdir "{temp_dir}" "{ppt_path}"'
                    logger.debug(f"Windows ODP转换命令: {cmd_str_odp}")
                    
                    # Windows下使用shell=True执行命令字符串
                    result_odp = subprocess.run(
                        cmd_str_odp,
                        shell=True,
                        capture_output=True,
                        text=True,
                        timeout=120,  # 2分钟超时
                        cwd=temp_dir  # 设置工作目录
                    )
                else:
                    # Linux/macOS下使用列表形式
                    logger.debug(f"执行ODP转换命令: {' '.join(cmd_to_odp)}")
                    result_odp = subprocess.run(
                        cmd_to_odp,
                        capture_output=True,
                        text=True,
                        timeout=120,
                        cwd=temp_dir
                    )
                
                # 检查ODP转换是否成功
                if result_odp.returncode != 0:
                    logger.error(f"转换为ODP失败: {result_odp.stderr}")
                    return False
                
                # 检查ODP文件是否生成
                odp_file = os.path.join(temp_dir, f"{ppt_name}.odp")
                if not os.path.exists(odp_file):
                    logger.error("未能生成ODP文件")
                    return False
                
                logger.info(f"成功转换为ODP格式: {odp_file}")

                # 步骤2.2: 从ODP转换为PDF格式
                cmd_to_pdf = [
                    self.libreoffice_cmd,
                    '--headless',           # 无头模式
                    '--invisible',          # 不可见
                    '--nodefault',          # 不使用默认设置
                    '--nolockcheck',        # 不检查锁定
                    '--nologo',             # 不显示logo
                    '--norestore',          # 不恢复
                    '--convert-to', 'pdf:writer_pdf_Export',  # 转换为PDF并指定导出过滤器
                    '--outdir', temp_dir,   # 输出目录
                    odp_file                # 输入文件
                ]

                # 在Windows下，如果路径包含空格，需要特殊处理
                if platform.system() == "Windows":
                    # 对于Windows，确保路径被正确引用
                    cmd_str_pdf = f'"{self.libreoffice_cmd}" --headless --invisible --nodefault --nolockcheck --nologo --norestore --convert-to pdf:writer_pdf_Export --outdir "{temp_dir}" "{odp_file}"'
                    logger.debug(f"Windows PDF转换命令: {cmd_str_pdf}")
                    
                    # 执行转换命令
                    start_time = time.time()
                    
                    # Windows下使用shell=True执行命令字符串
                    result = subprocess.run(
                        cmd_str_pdf,
                        shell=True,
                        capture_output=True,
                        text=True,
                        timeout=120,  # 2分钟超时
                        cwd=temp_dir  # 设置工作目录
                    )
                else:
                    # Linux/macOS下使用列表形式
                    logger.debug(f"执行PDF转换命令: {' '.join(cmd_to_pdf)}")
                    
                    # 执行转换命令
                    start_time = time.time()
                    
                    result = subprocess.run(
                        cmd_to_pdf,
                        capture_output=True,
                        text=True,
                        timeout=120,
                        cwd=temp_dir
                    )

                end_time = time.time()
                conversion_time = end_time - start_time

                logger.debug(f"LibreOffice转换耗时: {conversion_time:.2f}秒")
                logger.debug(f"返回码: {result.returncode}")

                if result.stdout:
                    logger.debug(f"标准输出: {result.stdout}")
                if result.stderr:
                    logger.debug(f"标准错误: {result.stderr}")

                # 检查转换是否成功
                if result.returncode == 0:
                    # 查找生成的PDF文件
                    try:
                        pdf_file = os.path.join(temp_dir, f"{ppt_name}.pdf")
                        if os.path.exists(pdf_file):
                            pdf_size = os.path.getsize(pdf_file)

                            logger.info(f"✅ PDF转换成功: {os.path.basename(pdf_file)} ({pdf_size} bytes)")
                            logger.info("🎯 PPT已被LibreOffice完整渲染，文本框自适应设置已生效")

                            self.stats['render_triggered'] = 1
                            self.stats['pdf_generated'] = 1
                            self.stats['pdf_deleted'] = 1  # PDF在临时目录中会自动删除

                            # 步骤2.3: 从ODP转换回PPTX格式
                            logger.info("步骤2.3: 从ODP转换回PPTX格式并覆盖原文件")
                            pptx_result = self._convert_odp_to_pptx(odp_file, temp_dir, ppt_path)
                            if pptx_result:
                                logger.info("✅ 完整渲染流程成功: ODP -> PDF -> PPTX")
                                self.stats['pptx_generated'] = 1
                            else:
                                logger.warning("⚠️ PDF渲染成功，但PPTX转换失败")

                            return True
                        else:
                            logger.warning("PDF转换命令成功但未找到PDF文件")
                            logger.debug(f"临时目录内容: {os.listdir(temp_dir)}")
                            return False
                    except Exception as e:
                        logger.error(f"检查PDF文件时出错: {e}")
                        return False
                else:
                    logger.error(f"LibreOffice转换失败，返回码: {result.returncode}")
                    if result.stderr:
                        logger.error(f"错误信息: {result.stderr}")
                    if result.stdout:
                        logger.error(f"输出信息: {result.stdout}")
                    return False

        except subprocess.TimeoutExpired:
            logger.error("LibreOffice转换超时（120秒）")
            return False
        except Exception as e:
            logger.error(f"PDF转换触发渲染失败: {e}")
            import traceback
            logger.debug(f"详细错误: {traceback.format_exc()}")
            return False

    def _convert_odp_to_pptx(self, odp_file: str, output_dir: str, original_pptx_path: str = None) -> str:
        """
        将ODP文件转换回PPTX格式
        
        Args:
            odp_file: ODP文件路径
            output_dir: 输出目录
            original_pptx_path: 原始PPTX文件路径，如果提供则会覆盖该文件
            
        Returns:
            str: 转换后的PPTX文件路径，失败则返回空字符串
        """
        try:
            logger.info(f"将ODP转换回PPTX格式: {os.path.basename(odp_file)}")
            
            # 构建转换命令
            cmd = [
                self.libreoffice_cmd,
                '--headless',
                '--invisible',
                '--nodefault',
                '--nolockcheck',
                '--nologo',
                '--norestore',
                '--convert-to', 'pptx:Impress MS PowerPoint 2007 XML',  # 显式指定导出格式
                '--outdir', output_dir,
                odp_file
            ]
            
            # 在Windows下，如果路径包含空格，需要特殊处理
            if platform.system() == "Windows":
                cmd_str = f'"{self.libreoffice_cmd}" --headless --invisible --nodefault --nolockcheck --nologo --norestore --convert-to "pptx:Impress MS PowerPoint 2007 XML" --outdir "{output_dir}" "{odp_file}"'
                logger.debug(f"Windows PPTX转换命令: {cmd_str}")
                
                # Windows下使用shell=True执行命令字符串
                result = subprocess.run(
                    cmd_str,
                    shell=True,
                    capture_output=True,
                    text=True,
                    timeout=120,
                    cwd=output_dir
                )
            else:
                # Linux/macOS下使用列表形式
                logger.debug(f"执行PPTX转换命令: {' '.join(cmd)}")
                result = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=120,
                    cwd=output_dir
                )
            
            if result.returncode == 0:
                # 查找生成的PPTX文件
                ppt_name = os.path.splitext(os.path.basename(odp_file))[0]
                pptx_file = os.path.join(output_dir, f"{ppt_name}.pptx")
                
                if os.path.exists(pptx_file):
                    logger.info(f"ODP成功转换为PPTX: {pptx_file}")
                    
                    # 如果提供了原始PPTX路径，则覆盖它
                    if original_pptx_path:
                        try:
                            shutil.copyfile(pptx_file, original_pptx_path)
                            logger.info(f"渲染后的PPTX已覆盖原文件: {original_pptx_path}")
                            return original_pptx_path
                        except Exception as e:
                            logger.error(f"无法覆盖原PPTX文件: {e}")
                            return pptx_file
                    
                    return pptx_file
                else:
                    logger.warning("PPTX转换命令成功但未找到PPTX文件")
                    logger.debug(f"目录内容: {os.listdir(output_dir)}")
                    return ""
            else:
                logger.error(f"PPTX转换失败，返回码: {result.returncode}")
                if result.stderr:
                    logger.error(f"错误信息: {result.stderr}")
                return ""
                
        except Exception as e:
            logger.error(f"ODP转换为PPTX出错: {e}")
            import traceback
            logger.debug(f"详细错误: {traceback.format_exc()}")
            return ""

    def _log_stats(self):
        """输出统计信息"""
        logger.info("=" * 50)
        logger.info("LibreOffice渲染触发统计")
        logger.info("=" * 50)
        logger.info(f"总文本框数: {self.stats['total_textboxes']}")
        logger.info(f"设置自适应: {self.stats['autofit_set']}")
        logger.info(f"渲染触发: {self.stats['render_triggered']}")
        logger.info(f"PDF生成: {self.stats['pdf_generated']}")
        logger.info(f"PDF删除: {self.stats['pdf_deleted']}")

        if self.stats['total_textboxes'] > 0:
            success_rate = (self.stats['autofit_set'] / self.stats['total_textboxes']) * 100
            logger.info(f"处理成功率: {success_rate:.1f}%")

    def get_stats(self) -> Dict[str, Any]:
        """获取统计信息"""
        return {
            'total_textboxes': self.stats.get('total_textboxes', 0),
            'autofit_set': self.stats.get('autofit_set', 0),
            'render_triggered': self.stats.get('render_triggered', 0),
            'pdf_generated': self.stats.get('pdf_generated', 0),
            'pdf_deleted': self.stats.get('pdf_deleted', 0),
            'pptx_generated': self.stats.get('pptx_generated', 0),
            'libreoffice_available': self.libreoffice_available,
            'libreoffice_cmd': getattr(self, 'libreoffice_cmd', None)
        }


def libreoffice_trigger_ppt_autofit(ppt_path: str) -> bool:
    """
    使用LibreOffice触发PPT文本框自适应渲染

    Args:
        ppt_path: PPT文件路径

    Returns:
        bool: 处理是否成功
    """
    trigger = LibreOfficeRenderTrigger()
    return trigger.process_ppt_with_render_trigger(ppt_path)


def install_libreoffice_instructions():
    """LibreOffice安装说明"""
    instructions = """
    LibreOffice安装说明:

    Ubuntu/Debian:
        sudo apt-get update
        sudo apt-get install -y libreoffice --no-install-recommends

    CentOS/RHEL:
        sudo yum install -y libreoffice-headless
        # 或者
        sudo dnf install -y libreoffice-headless

    Alpine Linux:
        apk add --no-cache libreoffice

    Docker Dockerfile示例:
        FROM python:3.9-slim
        RUN apt-get update && \\
            apt-get install -y libreoffice --no-install-recommends && \\
            rm -rf /var/lib/apt/lists/*

    验证安装:
        libreoffice --version
        # 或者
        soffice --version

    注意事项:
        1. 确保有足够的磁盘空间用于临时PDF文件
        2. 转换过程可能需要1-2分钟，取决于PPT大小
        3. 无头模式不需要图形界面
        4. 转换完成后PDF文件会自动删除
    """
    return instructions


# 测试函数
def test_libreoffice_availability():
    """测试LibreOffice可用性"""
    trigger = LibreOfficeRenderTrigger()

    if trigger.libreoffice_available:
        print(f"✅ LibreOffice可用: {trigger.libreoffice_cmd}")
        return True
    else:
        print("❌ LibreOffice不可用")
        print("\n安装说明:")
        print(install_libreoffice_instructions())
        return False


if __name__ == "__main__":
    # 测试LibreOffice可用性
    test_libreoffice_availability()
